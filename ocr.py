"""
High-Performance OCR Module using Docling
Supports: PDF, DOCX, PPTX, XLSX, and Images
RAG-friendly output with semantic chunking and rich metadata
"""

from __future__ import annotations

import base64
import hashlib
import io
import logging
import os
import re
import tempfile
import time
import warnings
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from typing import Any, Iterator, Optional

import ollama
try:
    from openai import OpenAI
    HAS_OPENAI = True
except ImportError:
    HAS_OPENAI = False
from PIL import Image

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from config import (
    VISION_MODEL, 
    VISION_MODEL_PRIMARY, 
    VISION_MODEL_SECONDARY,
    VISION_MODEL_LOCAL, 
    DATA_PATH
)

# Suppress deprecation warnings from docling_core
warnings.filterwarnings("ignore", category=DeprecationWarning, module="docling_core")

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Reduce noise from RapidOCR empty results (expected for non-text pages)
logging.getLogger("docling.models.rapid_ocr_model").setLevel(logging.ERROR)


# ============================================================================
# Configuration and Constants
# ============================================================================

class DocumentType(Enum):
    PDF = "pdf"
    DOCX = "docx"
    PPTX = "pptx"
    XLSX = "xlsx"
    IMAGE = "image"
    UNKNOWN = "unknown"


# Supported file extensions
SUPPORTED_EXTENSIONS = {
    ".pdf": DocumentType.PDF,
    ".docx": DocumentType.DOCX,
    ".doc": DocumentType.DOCX,
    ".pptx": DocumentType.PPTX,
    ".ppt": DocumentType.PPTX,
    ".xlsx": DocumentType.XLSX,
    ".xls": DocumentType.XLSX,
    ".png": DocumentType.IMAGE,
    ".jpg": DocumentType.IMAGE,
    ".jpeg": DocumentType.IMAGE,
    ".tiff": DocumentType.IMAGE,
    ".tif": DocumentType.IMAGE,
    ".bmp": DocumentType.IMAGE,
    ".webp": DocumentType.IMAGE,
}

# Constants VISION_MODEL and DATA_PATH imported from config


# ============================================================================
# Universal Metadata Extraction (Access Control)
# ============================================================================

def extract_file_metadata(file_path: Path) -> dict:
    """
    Extract universal metadata from file:
    - Access Level [L1]-[L5] from filename (Default: L2 = Internal)
    - Creation timestamp
    - Keywords from filename
    
    Access Level Guide:
    - L1: Public (blogs, brochures)
    - L2: Internal (default, standard docs)
    - L3: Confidential (project specifics, client data)
    - L4: Sensitive (HR, salaries, strategy)
    - L5: Top Secret (admin/owner only)
    """
    file_name = file_path.name
    
    # 1. Access Level (Default = 2: Internal)
    access_level = 2
    match = re.search(r"\[L([1-5])\]", file_name, re.IGNORECASE)
    if match:
        access_level = int(match.group(1))

    # 2. Creation timestamp
    try:
        created_at = os.path.getctime(file_path)
    except Exception:
        created_at = time.time()

    # 3. Keywords from filename
    clean_name = re.sub(r"\[L[1-5]\]", "", file_path.stem)
    keywords = [w.lower() for w in re.split(r"[_\-\s]+", clean_name) if len(w) > 2]
    # Remove common junk words
    junk = {"the", "and", "for", "with", "doc", "file", "pdf", "docx", "xlsx", "pptx"}
    keywords = [k for k in keywords if k not in junk]

    return {
        "access_level": access_level,
        "created_at": created_at,
        "keywords": ", ".join(keywords),
    }


# ============================================================================
# Data Classes for Structured Output
# ============================================================================

@dataclass
class ExtractedContent:
    """Structured container for extracted document content"""
    file_name: str
    file_path: str
    doc_type: DocumentType
    text_content: str
    tables: list[dict] = field(default_factory=list)
    images: list[dict] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)
    ocr_content: str = ""
    image_caption: str = ""
    page_count: int = 0
    word_count: int = 0
    # New: Heading context for better RAG
    headings: list[dict] = field(default_factory=list)  # List of {level, text, page}
    
    def __post_init__(self):
        self.word_count = len(self.text_content.split())
        self.metadata["file_name"] = self.file_name
        self.metadata["doc_type"] = self.doc_type.value
        self.metadata["word_count"] = self.word_count
        self.metadata["page_count"] = self.page_count


@dataclass
class ProcessingResult:
    """Result of document processing"""
    success: bool
    content: Optional[ExtractedContent] = None
    error: Optional[str] = None
    processing_time: float = 0.0


# ============================================================================
# OCR Engine Selection
# ============================================================================

class OCREngineType(Enum):
    """OCR Engine types"""
    TESSERACT = "tesseract"  # Best for clean documents, PDFs
    EASYOCR = "easyocr"      # Best for complex layouts, images, handwriting
    AUTO = "auto"            # Auto-select based on document type


def get_ocr_engine_for_type(doc_type: DocumentType) -> OCREngineType:
    """
    Select optimal OCR engine based on document type
    - Tesseract: Clean PDFs, Office docs (faster, good for structured text)
    - EasyOCR: Images, scanned docs (better for complex layouts)
    """
    if doc_type == DocumentType.IMAGE:
        return OCREngineType.EASYOCR
    elif doc_type in (DocumentType.PDF,):
        # PDFs can be mixed - use Tesseract for text-based, but Docling handles this
        return OCREngineType.TESSERACT
    else:
        # Office docs (DOCX, PPTX, XLSX) - Tesseract for any embedded images
        return OCREngineType.TESSERACT


# ============================================================================
# Docling Document Converter
# ============================================================================

class DoclingOCREngine:
    """
    High-performance document conversion engine using Docling
    Optimized for RAG pipelines with semantic chunking
    Uses Tesseract for documents and EasyOCR for images
    """
    
    def __init__(
        self,
        enable_ocr: bool = True,
        enable_tables: bool = True,
        max_workers: int = 4,
        vision_model: str = VISION_MODEL,
        ocr_engine: OCREngineType = OCREngineType.AUTO,
    ):
        self.enable_ocr = enable_ocr
        self.enable_tables = enable_tables
        self.max_workers = max_workers
        self.vision_model = vision_model
        self.ocr_engine = ocr_engine
        self._converter_tesseract = None
        self._converter_easyocr = None
        self._initialized_tesseract = False
        self._initialized_easyocr = False
    
    def _init_tesseract_converter(self):
        """Initialize Docling with Tesseract OCR for documents"""
        if self._initialized_tesseract:
            return self._converter_tesseract
        
        try:
            from docling.document_converter import DocumentConverter
            from docling.datamodel.pipeline_options import (
                PdfPipelineOptions,
                TableFormerMode,
                TesseractOcrOptions,
            )
            from docling.datamodel.base_models import InputFormat
            from docling.document_converter import PdfFormatOption
            
            pipeline_options = PdfPipelineOptions()
            pipeline_options.do_ocr = self.enable_ocr
            pipeline_options.do_table_structure = self.enable_tables
            
            if self.enable_tables:
                pipeline_options.table_structure_options.mode = TableFormerMode.ACCURATE
            
            # Use Tesseract for clean document OCR
            pipeline_options.ocr_options = TesseractOcrOptions(lang=["eng"])
            logger.info("Initialized Tesseract OCR engine for documents")
            
            self._converter_tesseract = DocumentConverter(
                format_options={
                    InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options),
                }
            )
            self._initialized_tesseract = True
            return self._converter_tesseract
            
        except ImportError as e:
            logger.warning(f"Tesseract not available: {e}")
            return None
    
    def _init_easyocr_converter(self):
        """Initialize Docling with EasyOCR for images and complex layouts"""
        if self._initialized_easyocr:
            return self._converter_easyocr
        
        try:
            from docling.document_converter import DocumentConverter
            from docling.datamodel.pipeline_options import (
                PdfPipelineOptions,
                TableFormerMode,
                EasyOcrOptions,
            )
            from docling.datamodel.base_models import InputFormat
            from docling.document_converter import PdfFormatOption
            
            pipeline_options = PdfPipelineOptions()
            pipeline_options.do_ocr = self.enable_ocr
            pipeline_options.do_table_structure = self.enable_tables
            
            if self.enable_tables:
                pipeline_options.table_structure_options.mode = TableFormerMode.ACCURATE
            
            # Use EasyOCR for images and complex layouts
            pipeline_options.ocr_options = EasyOcrOptions(lang=["en"])
            logger.info("Initialized EasyOCR engine for images")
            
            self._converter_easyocr = DocumentConverter(
                format_options={
                    InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options),
                }
            )
            self._initialized_easyocr = True
            return self._converter_easyocr
            
        except ImportError as e:
            logger.warning(f"EasyOCR not available: {e}")
            return None
    
    def _get_converter(self, doc_type: DocumentType):
        """Get appropriate converter based on document type"""
        if self.ocr_engine == OCREngineType.TESSERACT:
            return self._init_tesseract_converter()
        elif self.ocr_engine == OCREngineType.EASYOCR:
            return self._init_easyocr_converter()
        else:
            # AUTO mode - select based on document type
            engine = get_ocr_engine_for_type(doc_type)
            if engine == OCREngineType.EASYOCR:
                converter = self._init_easyocr_converter()
                if converter:
                    return converter
                # Fallback to Tesseract
                return self._init_tesseract_converter()
            else:
                converter = self._init_tesseract_converter()
                if converter:
                    return converter
                # Fallback to EasyOCR
                return self._init_easyocr_converter()
    
    def _get_document_type(self, file_path: Path) -> DocumentType:
        """Determine document type from file extension"""
        ext = file_path.suffix.lower()
        return SUPPORTED_EXTENSIONS.get(ext, DocumentType.UNKNOWN)
    
    def _compute_file_hash(self, file_path: Path) -> str:
        """Compute MD5 hash for file deduplication"""
        hasher = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                hasher.update(chunk)
        return hasher.hexdigest()
    
    def _encode_image_base64(self, image_path: Path) -> str:
        """Encode image to base64 for vision model"""
        with open(image_path, "rb") as f:
            return base64.b64encode(f.read()).decode("utf-8")
    
    def _generate_image_caption(self, image_path: Path) -> str:
        """
        Generate caption using vision model (Gemini -> Local Fallback)
        Returns: caption
        """
        caption_prompt = """Provide a detailed description of this image:
1. Document/Image Type: What kind of document or image is this?
2. Main Subject: What is the primary content or subject?
3. Key Elements: List important visual elements, diagrams, or figures.
4. Context: What is the likely purpose or use of this document?
5. Notable Details: Any important details that would help understand the content.

Be concise but comprehensive. Focus on information useful for search and retrieval."""
        
        caption = ""
        
        # 1. Try OpenRouter (Primary -> Secondary) if configured
        api_key = os.environ.get("OPENROUTER_API_KEY")
        if HAS_OPENAI and api_key:
            client = OpenAI(
                base_url="https://openrouter.ai/api/v1",
                api_key=api_key,
            )
            
            image_base64 = self._encode_image_base64(image_path)
            mime_type = "image/png" if image_path.suffix.lower() == ".png" else "image/jpeg"
            
            for model_name in [VISION_MODEL_PRIMARY, VISION_MODEL_SECONDARY]:
                try:
                    logger.info(f"Trying caption for {image_path.name} with OpenRouter model: {model_name}")
                    response_caption = client.chat.completions.create(
                        model=model_name,
                        messages=[
                            {
                                "role": "user",
                                "content": [
                                    {"type": "text", "text": caption_prompt},
                                    {
                                        "type": "image_url",
                                        "image_url": {
                                            "url": f"data:{mime_type};base64,{image_base64}"
                                        }
                                    }
                                ]
                            }
                        ]
                    )
                    
                    if response_caption.choices[0].message.content:
                        caption = response_caption.choices[0].message.content.strip()
                        logger.info(f"✓ Generated caption for {image_path.name} with OpenRouter ({model_name})")
                        return caption

                except Exception as e:
                    logger.warning(f"⚠ OpenRouter model {model_name} failed for {image_path.name}: {e}")
                    continue # Try next model
            
            logger.warning(f"⚠ All OpenRouter models failed. Falling back to Ollama.")

        # 2. Fallback to Local/Ollama
        try:
            image_base64 = self._encode_image_base64(image_path)
            
            # Generate caption
            caption_response = ollama.chat(
                model=VISION_MODEL_LOCAL,
                messages=[{
                    "role": "user",
                    "content": caption_prompt,
                    "images": [image_base64]
                }]
            )
            caption = caption_response["message"]["content"].strip()
            
        except Exception as e:
            logger.error(f"Error generating caption for {image_path.name}: {e}")
            caption = ""
        
        return caption
    
    def _process_with_docling(self, file_path: Path) -> ExtractedContent:
        """Process document using Docling with appropriate OCR engine"""
        doc_type = self._get_document_type(file_path)
        
        # Get appropriate converter (Tesseract for docs, EasyOCR for images)
        converter = self._get_converter(doc_type)
        
        try:
            if converter is None:
                raise ImportError("No Docling converter available")
            
            # Log which OCR engine is being used
            engine_name = get_ocr_engine_for_type(doc_type).value
            logger.info(f"Processing {file_path.name} with {engine_name} OCR")
            
            # Convert document
            result = converter.convert(str(file_path))
            doc = result.document
            
            # Extract text content with markdown formatting
            text_content = doc.export_to_markdown()
            
            # Extract tables with proper doc argument to avoid deprecation warning
            tables = []
            if hasattr(doc, 'tables'):
                for i, table in enumerate(doc.tables):
                    try:
                        # Pass doc argument to avoid deprecation warning
                        if hasattr(table, 'export_to_markdown'):
                            table_content = table.export_to_markdown(doc=doc)
                        else:
                            table_content = str(table)
                        table_data = {
                            "index": i,
                            "content": table_content,
                        }
                        tables.append(table_data)
                    except Exception as e:
                        logger.warning(f"Error extracting table {i}: {e}")
                        tables.append({"index": i, "content": str(table)})
            
            # Extract metadata
            metadata = {
                "file_hash": self._compute_file_hash(file_path),
                "converter": "docling",
                "ocr_engine": engine_name,
            }
            
            if hasattr(doc, 'num_pages'):
                page_count = doc.num_pages
            else:
                page_count = 1
            
            return ExtractedContent(
                file_name=file_path.name,
                file_path=str(file_path),
                doc_type=doc_type,
                text_content=text_content,
                tables=tables,
                metadata=metadata,
                page_count=page_count,
            )
            
        except Exception as e:
            logger.warning(f"Docling processing failed for {file_path.name}: {e}, trying fallback")
            return self._process_with_fallback(file_path, doc_type)
    
    def _process_with_fallback(self, file_path: Path, doc_type: DocumentType) -> ExtractedContent:
        """Fallback processing when Docling fails"""
        text_content = ""
        tables = []
        page_count = 1
        
        try:
            if doc_type == DocumentType.PDF:
                text_content, page_count = self._extract_pdf_fallback(file_path)
            elif doc_type == DocumentType.DOCX:
                text_content = self._extract_docx_fallback(file_path)
            elif doc_type == DocumentType.PPTX:
                text_content = self._extract_pptx_fallback(file_path)
            elif doc_type == DocumentType.XLSX:
                text_content, tables = self._extract_xlsx_fallback(file_path)
            else:
                text_content = f"Unsupported document type: {doc_type.value}"
                
        except Exception as e:
            logger.error(f"Fallback extraction failed for {file_path.name}: {e}")
            text_content = f"Failed to extract content from {file_path.name}"
        
        return ExtractedContent(
            file_name=file_path.name,
            file_path=str(file_path),
            doc_type=doc_type,
            text_content=text_content,
            tables=tables,
            metadata={"file_hash": self._compute_file_hash(file_path), "converter": "fallback"},
            page_count=page_count,
        )
    
    def _extract_pdf_fallback(self, file_path: Path) -> tuple[str, int]:
        """Fallback PDF extraction using pypdf"""
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(file_path))
            text_parts = []
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                text_parts.append(f"[Page {i + 1}]\n{page_text}")
            return "\n\n".join(text_parts), len(reader.pages)
        except Exception as e:
            logger.error(f"PDF fallback failed: {e}")
            return "", 0
    
    def _extract_docx_fallback(self, file_path: Path) -> str:
        """Fallback DOCX extraction using python-docx"""
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(str(file_path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception as e:
            logger.error(f"DOCX fallback failed: {e}")
            return ""
    
    def _extract_pptx_fallback(self, file_path: Path) -> str:
        """Fallback PPTX extraction using python-pptx"""
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            text_parts = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                slide_text.append(f"[Slide {i + 1}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                text_parts.append("\n".join(slide_text))
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"PPTX fallback failed: {e}")
            return ""
    
    def _extract_xlsx_fallback(self, file_path: Path) -> tuple[str, list[dict]]:
        """Fallback XLSX extraction using openpyxl"""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(file_path), data_only=True)
            text_parts = []
            tables = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"[Sheet: {sheet_name}]")
                
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_values):
                        rows.append(" | ".join(row_values))
                
                sheet_content = "\n".join(rows)
                text_parts.append(sheet_content)
                tables.append({"sheet": sheet_name, "content": sheet_content})
            
            return "\n\n".join(text_parts), tables
        except Exception as e:
            logger.error(f"XLSX fallback failed: {e}")
            return "", []
    
    def _process_image(self, file_path: Path) -> ExtractedContent:
        """Process image file: Classical OCR + Vision Caption"""
        
        # 1. Classical OCR (using Docling/EasyOCR/Tesseract)
        # We reuse _process_with_docling because it handles image OCR via EasyOCR
        try:
            content = self._process_with_docling(file_path)
            ocr_text = content.text_content
        except Exception as e:
            logger.warning(f"Standard OCR failed for {file_path.name}: {e}")
            ocr_text = ""
            content = ExtractedContent(
                file_name=file_path.name,
                file_path=str(file_path),
                doc_type=DocumentType.IMAGE,
                text_content="",
                metadata={"file_hash": self._compute_file_hash(file_path), "converter": "vision_fallback"}
            )

        # 2. Semantic Caption (Vision Model)
        caption = self._generate_image_caption(file_path)
        
        # 3. Combine Content
        content_parts = []
        content_parts.append(f"# Image: {file_path.name}")
        content_parts.append("")
        
        if caption:
            content_parts.append("## Image Description")
            content_parts.append(caption)
            content_parts.append("")
        
        if ocr_text:
            content_parts.append("## Extracted Text (OCR)")
            content_parts.append(ocr_text)
        
        content.text_content = "\n".join(content_parts)
        content.image_caption = caption
        content.metadata["has_caption"] = bool(caption)
        content.metadata["has_ocr"] = bool(ocr_text)
        content.metadata["converter"] = "docling_ocr+vision_caption"
        
        return content
    
    def process_file(self, file_path: Path | str) -> ProcessingResult:
        """
        Process a single file and extract content
        
        Args:
            file_path: Path to the file to process
            
        Returns:
            ProcessingResult with extracted content or error
        """
        import time
        start_time = time.time()
        
        file_path = Path(file_path)
        
        if not file_path.exists():
            return ProcessingResult(
                success=False,
                error=f"File not found: {file_path}"
            )
        
        doc_type = self._get_document_type(file_path)
        
        if doc_type == DocumentType.UNKNOWN:
            return ProcessingResult(
                success=False,
                error=f"Unsupported file type: {file_path.suffix}"
            )
        
        try:
            if doc_type == DocumentType.IMAGE:
                content = self._process_image(file_path)
            else:
                content = self._process_with_docling(file_path)
            
            processing_time = time.time() - start_time
            
            return ProcessingResult(
                success=True,
                content=content,
                processing_time=processing_time
            )
            
        except Exception as e:
            return ProcessingResult(
                success=False,
                error=str(e),
                processing_time=time.time() - start_time
            )
    
    def process_directory(
        self,
        directory: Path | str,
        recursive: bool = False,
    ) -> Iterator[ProcessingResult]:
        """
        Process all supported files in a directory
        
        Args:
            directory: Path to directory
            recursive: Whether to process subdirectories
            
        Yields:
            ProcessingResult for each file
        """
        directory = Path(directory)
        
        if not directory.exists():
            logger.error(f"Directory not found: {directory}")
            return
        
        # Collect all supported files
        pattern = "**/*" if recursive else "*"
        files = [
            f for f in directory.glob(pattern)
            if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS
        ]
        
        logger.info(f"Found {len(files)} files to process in {directory}")
        
        # Process files with thread pool for parallel processing
        with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
            future_to_file = {
                executor.submit(self.process_file, f): f for f in files
            }
            
            for future in as_completed(future_to_file):
                file_path = future_to_file[future]
                try:
                    result = future.result()
                    yield result
                except Exception as e:
                    yield ProcessingResult(
                        success=False,
                        error=f"Processing error for {file_path}: {e}"
                    )


# ============================================================================
# RAG-Friendly Document Processor
# ============================================================================

class RAGDocumentProcessor:
    """
    Converts extracted content to RAG-friendly LangChain Documents
    with semantic chunking, heading context, and access control metadata
    """
    
    def __init__(
        self,
        chunk_size: int = 512,
        chunk_overlap: int = 128,
        add_semantic_headers: bool = False,  # Disabled - metadata stored separately
    ):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.add_semantic_headers = add_semantic_headers
        
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=[
                "\n\n## ",  # Markdown H2 headers
                "\n\n### ", # Markdown H3 headers
                "\n\n",     # Double newline (paragraphs)
                "\n",       # Single newline
                ". ",       # Sentences
                " ",        # Words
                "",         # Characters
            ],
        )
        
        # Pattern to detect headings (numbered sections, uppercase lines, etc.)
        self.heading_patterns = [
            r"^(\d+\.)+\s+.+$",           # 1. or 1.1. or 1.1.1. style
            r"^[A-Z][A-Z\s]{5,50}$",      # ALL CAPS lines (5-50 chars)
            r"^(Chapter|Section|Part)\s+\d+",  # Chapter X, Section Y
            r"^#{1,3}\s+.+$",              # Markdown headers
        ]
        
        # Patterns to detect contact blocks
        self.contact_patterns = [
            r'(?:Dr\.|Mr\.|Ms\.|Sh\.|Smt\.)\s*[\w\s]+',  # Names with titles
            r'\d{3}[-.\s]?\d{4,8}',  # Phone numbers
            r'[\w\.]+\[(?:at|dot)\][\w\.]+',  # Obfuscated emails
            r'[\w\.\-]+@[\w\.\-]+',  # Regular emails
            r'(?:Room|Block|Floor)\s*(?:No\.?)?\s*[\d\-]+',  # Office locations
        ]
    
    def _detect_heading(self, line: str) -> bool:
        """Check if a line looks like a heading"""
        line = line.strip()
        if not line or len(line) > 100:  # Too long to be a heading
            return False
        for pattern in self.heading_patterns:
            if re.match(pattern, line):
                return True
        return False
    
    def _detect_contact_block(self, text: str) -> list[tuple[int, int, str]]:
        """
        Detect contact information blocks in text.
        Returns list of (start_idx, end_idx, block_content) tuples.
        """
        blocks = []
        lines = text.split('\n')
        
        i = 0
        while i < len(lines):
            line = lines[i].strip()
            
            # Check if line starts a contact entry (numbered person or title)
            if re.match(r'^\d+\.\s*(?:Dr\.|Mr\.|Ms\.|Sh\.|Smt\.)', line) or \
               re.match(r'^(?:Dr\.|Mr\.|Ms\.|Sh\.|Smt\.)\s+\w', line):
                
                # Found start of contact block, collect until next numbered entry or end
                block_start = sum(len(l) + 1 for l in lines[:i])
                block_lines = [lines[i]]
                j = i + 1
                
                # Collect all lines until next contact entry or empty section
                while j < len(lines):
                    next_line = lines[j].strip()
                    # Stop at next numbered person entry
                    if re.match(r'^\d+\.\s*(?:Dr\.|Mr\.|Ms\.|Sh\.|Smt\.)', next_line):
                        break
                    # Stop at major section changes
                    if next_line and self._detect_heading(next_line):
                        break
                    block_lines.append(lines[j])
                    j += 1
                
                block_content = '\n'.join(block_lines)
                block_end = block_start + len(block_content)
                
                # Only count as contact block if it has actual contact data
                has_phone = bool(re.search(r'\d{3}[-.\s]?\d{4,8}', block_content))
                has_email = bool(re.search(r'[\w\.]+(?:\[at\]|@)[\w\.]+', block_content))
                
                if has_phone or has_email:
                    blocks.append((block_start, block_end, block_content))
                
                i = j
            else:
                i += 1
        
        return blocks
    
    def _extract_contact_blocks(self, text: str) -> tuple[str, list[str]]:
        """
        Extract contact blocks from text, returning:
        - text_without_contacts: Original text with contact blocks replaced by markers
        - contact_blocks: List of contact block strings to be chunked separately
        """
        contact_blocks = []
        blocks = self._detect_contact_block(text)
        
        if not blocks:
            return text, []
        
        # Sort blocks by position (reverse to replace from end)
        blocks.sort(key=lambda x: x[0], reverse=True)
        
        text_modified = text
        for start, end, block in blocks:
            contact_blocks.append(block)
            # Replace with marker that will be split on
            marker = f"\n\n[CONTACT_BLOCK_{len(contact_blocks)}]\n\n"
            text_modified = text_modified[:start] + marker + text_modified[end:]
        
        # Reverse to maintain original order
        contact_blocks.reverse()
        
        return text_modified, contact_blocks
    
    def _create_chunk_id(self, file_name: str, chunk_index: int, page: int = 0) -> str:
        """Create unique chunk ID"""
        return f"{file_name}:p{page}:c{chunk_index}"
    
    def _add_semantic_header(self, content: str, metadata: dict) -> str:
        """Add semantic header to chunk for better retrieval"""
        if not self.add_semantic_headers:
            return content
        
        header_parts = []
        
        file_name = metadata.get("file_name", "unknown")
        doc_type = metadata.get("doc_type", "document")
        access_level = metadata.get("access_level", 2)
        
        header_parts.append(f"[Source: {file_name}]")
        header_parts.append(f"[Type: {doc_type}]")
        header_parts.append(f"[Access: L{access_level}]")
        
        if metadata.get("page_num"):
            header_parts.append(f"[Page: {metadata['page_num']}]")
        
        if metadata.get("sheet"):
            header_parts.append(f"[Sheet: {metadata['sheet']}]")
            
        # Add heading context if available
        if metadata.get("heading_context"):
            header_parts.append(f"[Section: {metadata['heading_context']}]")
        
        header = " ".join(header_parts)
        return f"{header}\n\n{content}"
    
    def process_extracted_content(self, content: ExtractedContent) -> list[Document]:
        """
        Convert ExtractedContent to list of LangChain Documents
        with access control and heading context
        
        Args:
            content: ExtractedContent from OCR engine
            
        Returns:
            List of chunked LangChain Documents
        """
        documents = []
        
        # Extract file-level metadata (access level, keywords, etc.)
        file_meta = extract_file_metadata(Path(content.file_path))
        
        # Create base metadata with simple types only
        base_metadata = {
            "source": content.file_path,
            "file_name": content.file_name,
            "doc_type": content.doc_type.value if hasattr(content.doc_type, 'value') else str(content.doc_type),
            "page_count": int(content.page_count() if callable(content.page_count) else content.page_count), 
            "word_count": int(content.word_count() if callable(content.word_count) else content.word_count),
            "access_level": int(file_meta.get("access_level", 2)),
            "created_at": float(file_meta.get("created_at", 0.0)),
            "keywords": str(file_meta.get("keywords", "")),
        }
        
        # Clean up any complex objects from content.metadata before merging
        for k, v in content.metadata.items():
            if isinstance(v, (str, int, float, bool)):
                base_metadata[k] = v
            else:
                 # Convert complex types (like DoclingDocument) to string or ignore
                 base_metadata[k] = str(v)
        
        # Special handling for images
        if content.doc_type == DocumentType.IMAGE:
            return self._process_image_content(content, base_metadata)
        
        # Track heading context for better RAG
        current_heading = "General"
        
        # Split text content into chunks with heading tracking
        if content.text_content.strip():
            # First pass: identify headings in the text
            lines = content.text_content.split('\n')
            text_with_heading_context = []
            
            for line in lines:
                if self._detect_heading(line):
                    current_heading = line.strip()[:80]  # Cap heading length
                text_with_heading_context.append((line, current_heading))
            
            # === Table-Aware Chunking: Extract contact blocks first ===
            text_for_chunking, contact_blocks = self._extract_contact_blocks(content.text_content)
            
            # Process contact blocks as separate chunks (keep together)
            contact_chunk_idx = 0
            for block in contact_blocks:
                if len(block.strip()) > 50:  # Only if meaningful content
                    contact_meta = base_metadata.copy()
                    contact_meta["content_type"] = "contact"
                    contact_meta["heading_context"] = "Contact Information"
                    contact_meta["chunk_id"] = self._create_chunk_id(
                        content.file_name, contact_chunk_idx, page=99  # Special page for contacts
                    )
                    contact_meta["chunk_index"] = contact_chunk_idx
                    contact_meta["is_contact_block"] = True
                    
                    contact_doc = Document(
                        page_content=self._add_semantic_header(block, contact_meta),
                        metadata=contact_meta
                    )
                    documents.append(contact_doc)
                    contact_chunk_idx += 1
                    logger.debug(f"Created contact block chunk: {block[:100]}...")
            
            # Now chunk the remaining text (with contact blocks removed)
            # Remove the [CONTACT_BLOCK_X] markers
            text_for_chunking = re.sub(r'\[CONTACT_BLOCK_\d+\]', '', text_for_chunking)
            
            if text_for_chunking.strip():
                text_doc = Document(
                    page_content=text_for_chunking,
                    metadata=base_metadata.copy()
                )
                
                chunks = self.text_splitter.split_documents([text_doc])
                
                # Offset chunk indices by number of contact blocks
                chunk_offset = len(contact_blocks)
                
                for i, chunk in enumerate(chunks):
                    # Try to find the heading context for this chunk
                    chunk_start = chunk.page_content[:100]
                    heading_for_chunk = "General"
                    
                    for line, heading in text_with_heading_context:
                        if line.strip() and line.strip() in chunk_start:
                            heading_for_chunk = heading
                            break
                    
                    chunk.metadata["chunk_id"] = self._create_chunk_id(
                        content.file_name, i + chunk_offset
                    )
                    chunk.metadata["chunk_index"] = i + chunk_offset
                    chunk.metadata["total_chunks"] = len(chunks) + len(contact_blocks)
                    chunk.metadata["heading_context"] = heading_for_chunk
                    chunk.metadata["content_type"] = "text"
                    
                    # Add semantic header with context
                    chunk.page_content = self._add_semantic_header(
                        chunk.page_content,
                        chunk.metadata
                    )
                    
                    documents.append(chunk)
        
        # Process tables as separate documents
        for table in content.tables:
            table_content = table.get("content", "")
            if table_content.strip():
                table_metadata = base_metadata.copy()
                table_metadata["content_type"] = "table"
                table_metadata["table_index"] = table.get("index", 0)
                table_metadata["heading_context"] = table.get("heading", "Table")
                
                if "sheet" in table:
                    table_metadata["sheet"] = table["sheet"]
                
                table_doc = Document(
                    page_content=self._add_semantic_header(table_content, table_metadata),
                    metadata=table_metadata
                )
                table_doc.metadata["chunk_id"] = self._create_chunk_id(
                    content.file_name,
                    len(documents),
                )
                documents.append(table_doc)
        
        return documents
    
    def _process_image_content(
        self,
        content: ExtractedContent,
        base_metadata: dict
    ) -> list[Document]:
        """Process image content with OCR, caption, and access control"""
        documents = []
        
        # Combine caption and OCR for comprehensive image representation
        image_metadata = base_metadata.copy()
        image_metadata["content_type"] = "image"
        image_metadata["has_ocr"] = bool(content.ocr_content)
        image_metadata["has_caption"] = bool(content.image_caption)
        image_metadata["heading_context"] = "Image Content"
        
        # Main image document with all content
        main_content = content.text_content
        
        if main_content.strip():
            chunks = self.text_splitter.split_documents([
                Document(page_content=main_content, metadata=image_metadata)
            ])
            
            for i, chunk in enumerate(chunks):
                chunk.metadata["chunk_id"] = self._create_chunk_id(
                    content.file_name, i
                )
                chunk.metadata["chunk_index"] = i
                # Add semantic header with access level
                chunk.page_content = self._add_semantic_header(
                    chunk.page_content,
                    chunk.metadata
                )
                documents.append(chunk)
        
        return documents


# ============================================================================
# Main Processing Functions (LangChain Integration)
# ============================================================================

def process_documents(
    data_path: str = DATA_PATH,
    chunk_size: int = 512,
    chunk_overlap: int = 128,
    max_workers: int = 4,
) -> list[Document]:
    """
    Process all supported documents in directory and return RAG-ready chunks
    
    This is the main entry point for integrating with LangChain pipelines.
    
    Args:
        data_path: Path to directory containing documents
        chunk_size: Size of text chunks
        chunk_overlap: Overlap between chunks
        max_workers: Number of parallel workers
        
    Returns:
        List of LangChain Documents ready for vector store
    """
    engine = DoclingOCREngine(
        enable_ocr=True,
        enable_tables=True,
        max_workers=max_workers,
    )
    
    processor = RAGDocumentProcessor(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
    )
    
    all_documents = []
    
    logger.info(f"Processing documents from: {data_path}")
    
    for result in engine.process_directory(data_path):
        if result.success and result.content:
            docs = processor.process_extracted_content(result.content)
            all_documents.extend(docs)
            logger.info(
                f"✓ Processed: {result.content.file_name} "
                f"({len(docs)} chunks, {result.processing_time:.2f}s)"
            )
        else:
            logger.warning(f"✗ Failed: {result.error}")
    
    logger.info(f"Total documents processed: {len(all_documents)} chunks")
    return all_documents


def process_specific_files(
    files: list[Path],
    chunk_size: int = 512,
    chunk_overlap: int = 128,
    max_workers: int = 4,
) -> list[Document]:
    """
    Process a specific list of files.
    """
    if not files:
        return []

    engine = DoclingOCREngine(
        enable_ocr=True,
        enable_tables=True,
        max_workers=max_workers,
    )
    
    processor = RAGDocumentProcessor(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
    )
    
    all_documents = []
    
    logger.info(f"Processing {len(files)} specific files")
    
    for file_path in files:
        if not file_path.exists():
            continue

        result = engine.process_file(file_path)
        if result.success and result.content:
            docs = processor.process_extracted_content(result.content)
            all_documents.extend(docs)
            logger.info(
                f"✓ Processed: {result.content.file_name} "
                f"({len(docs)} chunks, {result.processing_time:.2f}s)"
            )
        else:
            logger.warning(f"✗ Failed {file_path.name}: {result.error}")
            
    return all_documents


def process_pdfs_with_ocr(data_path: str = DATA_PATH) -> list[Document]:
    """Process only PDF files with OCR"""
    engine = DoclingOCREngine(enable_ocr=True, enable_tables=True)
    processor = RAGDocumentProcessor()
    
    documents = []
    data_dir = Path(data_path)
    
    for file_path in data_dir.glob("*.pdf"):
        result = engine.process_file(file_path)
        if result.success and result.content:
            docs = processor.process_extracted_content(result.content)
            documents.extend(docs)
            print(f"✓ PDF: {file_path.name} → {len(docs)} chunks")
    
    return documents


def process_office_documents(data_path: str = DATA_PATH) -> list[Document]:
    """Process DOCX, PPTX, XLSX files"""
    engine = DoclingOCREngine(enable_ocr=True, enable_tables=True)
    processor = RAGDocumentProcessor()
    
    documents = []
    data_dir = Path(data_path)
    
    office_extensions = [".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls"]
    
    for ext in office_extensions:
        for file_path in data_dir.glob(f"*{ext}"):
            result = engine.process_file(file_path)
            if result.success and result.content:
                docs = processor.process_extracted_content(result.content)
                documents.extend(docs)
                print(f"✓ Office: {file_path.name} → {len(docs)} chunks")
    
    return documents


def process_images_with_ocr_and_caption(data_path: str = DATA_PATH) -> list[Document]:
    """
    Process images with both OCR and vision model captioning
    Returns documents with combined OCR text and semantic description
    """
    engine = DoclingOCREngine()
    processor = RAGDocumentProcessor()
    
    documents = []
    data_dir = Path(data_path)
    
    image_extensions = [".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp"]
    
    for ext in image_extensions:
        for file_path in data_dir.glob(f"*{ext}"):
            result = engine.process_file(file_path)
            if result.success and result.content:
                docs = processor.process_extracted_content(result.content)
                documents.extend(docs)
                
                has_ocr = "✓" if result.content.ocr_content else "✗"
                print(f"✓ Image: {file_path.name} → {len(docs)} chunks (OCR: {has_ocr})")
    
    return documents


# ============================================================================
# Utility Functions
# ============================================================================

def get_supported_extensions() -> list[str]:
    """Return list of supported file extensions"""
    return list(SUPPORTED_EXTENSIONS.keys())


def is_supported_file(file_path: Path | str) -> bool:
    """Check if file type is supported"""
    path = Path(file_path)
    return path.suffix.lower() in SUPPORTED_EXTENSIONS


def get_file_stats(data_path: str = DATA_PATH) -> dict[str, int]:
    """Get statistics about files in directory"""
    data_dir = Path(data_path)
    stats = {"total": 0}
    
    for ext, doc_type in SUPPORTED_EXTENSIONS.items():
        count = len(list(data_dir.glob(f"*{ext}")))
        if count > 0:
            stats[ext] = count
            stats["total"] += count
    
    return stats


# ============================================================================
# Main Entry Point
# ============================================================================

if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description="Process documents with Docling OCR")
    parser.add_argument("--data-path", default=DATA_PATH, help="Path to documents")
    parser.add_argument("--type", choices=["all", "pdf", "office", "image"], default="all")
    args = parser.parse_args()
    
    print(f"\n{'='*60}")
    print("Docling OCR Document Processor")
    print(f"{'='*60}\n")
    
    # Show file statistics
    stats = get_file_stats(args.data_path)
    print(f"Files found: {stats}")
    print()
    
    # Process based on type
    if args.type == "all":
        documents = process_documents(args.data_path)
    elif args.type == "pdf":
        documents = process_pdfs_with_ocr(args.data_path)
    elif args.type == "office":
        documents = process_office_documents(args.data_path)
    elif args.type == "image":
        documents = process_images_with_ocr_and_caption(args.data_path)
    
    print(f"\n{'='*60}")
    print(f"Processing Complete: {len(documents)} total chunks")
    print(f"{'='*60}")
    
    # Show sample output
    if documents:
        print("\n--- Sample Chunks ---")
        for i, doc in enumerate(documents[:3]):
            print(f"\nChunk {i + 1}:")
            print(f"  ID: {doc.metadata.get('chunk_id', 'N/A')}")
            print(f"  Type: {doc.metadata.get('doc_type', 'N/A')}")
            print(f"  Source: {doc.metadata.get('file_name', 'N/A')}")
            print(f"  Content: {doc.page_content[:150]}...")
